# app/services/presentation_service.py
import re
import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt

# Define storage path
PRESENTATION_DIR = Path("app/storage/presentations")
PRESENTATION_DIR.mkdir(parents=True, exist_ok=True)

class PresentationService:
    """Service for creating well-structured and paginated PowerPoint presentations."""

    def _clean_text(self, text: str) -> str:
        """Removes markdown-like formatting for presentation."""
        text = re.sub(r'^\s*([\*\-]|\d+\.)\s*', '', text)
        text = re.sub(r'^\s*#+\s*', '', text)
        text = text.replace('**', '').replace('*', '')
        return text.strip()

    def _paginate_content(self, lines: list[str], max_words: int, max_chars: int) -> list[list[str]]:
        """
        Splits text content into multiple pages based on word and character limits,
        ensuring no page is overly dense.
        """
        pages = []
        # Join and clean all lines, then split into words to handle long paragraphs
        full_text = " ".join(self._clean_text(line) for line in lines)
        words = full_text.split()

        if not words:
            return []

        current_page_words = []
        for word in words:
            # Check if adding the next word would exceed the limits
            if current_page_words and (
                len(" ".join(current_page_words) + " " + word) > max_chars or
                len(current_page_words) + 1 > max_words
            ):
                # Finalize the current page and start a new one
                pages.append([" ".join(current_page_words)])
                current_page_words = [word]
            else:
                current_page_words.append(word)

        # Add the last remaining page
        if current_page_words:
            pages.append([" ".join(current_page_words)])

        return pages

    def create_presentation_from_content(self, document_id: str, json_content: str) -> Path:
        """
        Creates a .pptx presentation from a JSON structure containing the title and slides.
        """
        try:
            content = json.loads(json_content)
        except json.JSONDecodeError:
            raise ValueError("Invalid JSON content received for presentation.")

        prs = Presentation()
        
        # --- Title Slide ---
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1] if len(slide.placeholders) > 1 else None
        
        title.text = content.get("title", "Presentation")
        if subtitle:
            subtitle.text = f"Generated from document: {document_id}"

        # --- Content Slides ---
        content_slide_layout = prs.slide_layouts[1]
        for slide_data in content.get("slides", []):
            slide = prs.slides.add_slide(content_slide_layout)
            title_shape = slide.shapes.title
            body_shape = slide.placeholders[1]
            
            title_shape.text = slide_data.get("title", "")
            
            text_frame = body_shape.text_frame
            text_frame.clear()
            
            for point in slide_data.get("content", []):
                p = text_frame.add_paragraph()
                p.text = self._clean_text(point)
                p.level = 0
                p.font.size = Pt(18)

        # --- Save Presentation ---
        file_path = PRESENTATION_DIR / f"{document_id}.pptx"
        prs.save(file_path)
        
        return file_path

# Create a singleton instance
presentation_service = PresentationService()